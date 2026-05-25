from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────
BLACK    = RGBColor(0x09, 0x09, 0x0b)
WHITE    = RGBColor(0xfa, 0xfa, 0xf9)
OFFWHITE = RGBColor(0xf4, 0xf4, 0xf2)
DGRAY    = RGBColor(0x18, 0x18, 0x1b)
MGRAY    = RGBColor(0x27, 0x27, 0x2a)
BORDER   = RGBColor(0x3f, 0x3f, 0x46)
LGRAY    = RGBColor(0xa1, 0xa1, 0xaa)
MUTED    = RGBColor(0x71, 0x71, 0x7a)
GREEN    = RGBColor(0x22, 0xc5, 0x5e)
RED      = RGBColor(0xef, 0x44, 0x44)
AMBER    = RGBColor(0xf5, 0x9e, 0x0b)
BLUE     = RGBColor(0x3b, 0x82, 0xf6)
CYAN     = RGBColor(0x22, 0xd3, 0xee)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────
def slide(color=BLACK):
    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = color
    return sl

def rect(sl, l, t, w, h, fill=None, border=None, bw=Pt(0.75), rounding=False):
    kind = 9 if rounding else 1
    s = sl.shapes.add_shape(kind, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def tb(sl, text, l, t, w, h, size=Pt(14), bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True, spacing=None):
    tf = sl.shapes.add_textbox(l, t, w, h)
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf

def label_txt(sl, text, t=Inches(0.52), dark=True):
    c = MUTED if dark else RGBColor(0x99, 0x99, 0x99)
    tf = sl.shapes.add_textbox(Inches(0.75), t, Inches(6), Inches(0.28))
    p = tf.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(9)
    r.font.color.rgb = c
    r.font.bold = False

def rule_line(sl, t, dark=True):
    s = sl.shapes.add_shape(1, Inches(0.75), t, Inches(0.5), Pt(1))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE if dark else BLACK
    s.line.fill.background()

def h1(sl, text, t, w=Inches(11.8), size=Pt(42), dark=True, bold_parts=None):
    tc = WHITE if dark else BLACK
    tf = sl.shapes.add_textbox(Inches(0.75), t, w, Inches(2.0))
    tf.word_wrap = True
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if not bold_parts:
        r = p.add_run(); r.text = text
        r.font.size = size; r.font.bold = False; r.font.color.rgb = tc
    else:
        rem = text
        for part in bold_parts:
            idx = rem.find(part)
            if idx < 0: continue
            if idx > 0:
                r = p.add_run(); r.text = rem[:idx]
                r.font.size = size; r.font.bold = False; r.font.color.rgb = tc
            r = p.add_run(); r.text = part
            r.font.size = size; r.font.bold = True; r.font.color.rgb = tc
            rem = rem[idx+len(part):]
        if rem:
            r = p.add_run(); r.text = rem
            r.font.size = size; r.font.bold = False; r.font.color.rgb = tc

def item_list(sl, items, l, t, w, h, dark=True, num_col=None):
    body_c = WHITE if dark else BLACK
    num_c  = num_col or (MGRAY if dark else RGBColor(0xcc, 0xcc, 0xcc))
    sub_c  = LGRAY if dark else MUTED
    tf = sl.shapes.add_textbox(l, t, w, h)
    tf.word_wrap = True
    tf.text_frame.word_wrap = True
    first = True
    for i, item in enumerate(items, 1):
        p = tf.text_frame.paragraphs[0] if first else tf.text_frame.add_paragraph()
        first = False
        p.space_before = Pt(0) if i == 1 else Pt(14)
        p.space_after = Pt(0)
        r1 = p.add_run()
        r1.text = f"{i:02d}   "
        r1.font.size = Pt(10); r1.font.bold = True; r1.font.color.rgb = num_c
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(15); r2.font.color.rgb = body_c

def callout(sl, tag, body, l, t, w, h, dark=True):
    fill  = MGRAY if dark else OFFWHITE
    brd   = BORDER if dark else RGBColor(0xe0, 0xe0, 0xdd)
    tag_c = MUTED if dark else RGBColor(0xaa, 0xaa, 0xaa)
    body_c = LGRAY if dark else MUTED
    rect(sl, l, t, w, h, fill=fill, border=brd)
    tb(sl, tag.upper(), l+Inches(0.22), t+Inches(0.18), w-Inches(0.3), Inches(0.26),
       size=Pt(8), bold=True, color=tag_c)
    tb(sl, body, l+Inches(0.22), t+Inches(0.52), w-Inches(0.3), h-Inches(0.65),
       size=Pt(14), color=body_c, wrap=True)

def grid_cell_dark(sl, num, title, body, l, t, w, h):
    rect(sl, l, t, w, h, fill=MGRAY, border=BORDER)
    tb(sl, f"{num:02d}", l+Inches(0.24), t+Inches(0.2),  Inches(0.5),  Inches(0.28), size=Pt(9),  color=MUTED)
    tb(sl, title,       l+Inches(0.24), t+Inches(0.55),  w-Inches(0.4),Inches(0.38), size=Pt(14), bold=True, color=WHITE, wrap=True)
    tb(sl, body,        l+Inches(0.24), t+Inches(1.0),   w-Inches(0.4),h-Inches(1.2),size=Pt(12), color=LGRAY, wrap=True)

def grid_cell_light(sl, num, title, body, l, t, w, h):
    rect(sl, l, t, w, h, fill=WHITE, border=RGBColor(0xe0,0xe0,0xdd))
    tb(sl, f"{num:02d}", l+Inches(0.24), t+Inches(0.2),  Inches(0.5),  Inches(0.28), size=Pt(9),  color=RGBColor(0xcc,0xcc,0xcc))
    tb(sl, title,       l+Inches(0.24), t+Inches(0.55),  w-Inches(0.4),Inches(0.38), size=Pt(14), bold=True, color=BLACK, wrap=True)
    tb(sl, body,        l+Inches(0.24), t+Inches(1.0),   w-Inches(0.4),h-Inches(1.2),size=Pt(12), color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# S1 — TITLE
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)

# Thin top accent line
rect(sl, Inches(0), Inches(0), W, Pt(3), fill=RGBColor(0x22,0x22,0x26))

tb(sl, "VERKADA", Inches(0.75), Inches(0.45), Inches(3), Inches(0.32),
   size=Pt(11), bold=True, color=WHITE)

rule_line(sl, Inches(1.05))

# Large title
tf = sl.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(9.5), Inches(2.8))
tf.word_wrap = True
p = tf.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Zero-Trust Network Architecture\nfor Physical Security Infrastructure"
r.font.size = Pt(52); r.font.bold = False; r.font.color.rgb = WHITE

tb(sl, "A real scenario. A real threat. Built, deployed, and running in under two minutes.",
   Inches(0.75), Inches(4.15), Inches(8.5), Inches(0.65),
   size=Pt(18), color=LGRAY, wrap=True)

# Bottom bar
rect(sl, Inches(0), Inches(6.95), W, Inches(0.55), fill=MGRAY)
tb(sl, "SAMYUKTHA SYAMSUNDAR   ·   SOLUTIONS ENGINEER CANDIDATE   ·   VERKADA",
   Inches(0.75), Inches(7.0), Inches(11), Inches(0.38),
   size=Pt(10), color=MUTED)

# ═══════════════════════════════════════════════════════════════════
# S2 — ABOUT ME
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "About Me")
rule_line(sl, Inches(0.88))
h1(sl, "Samyuktha Syamsundar", Inches(0.98), size=Pt(46))
tb(sl, "B.S. Cybersecurity, Purdue University   ·   St. Louis, Missouri",
   Inches(0.75), Inches(2.28), Inches(10), Inches(0.38),
   size=Pt(17), color=LGRAY)

cards = [
    ("Reading",  "Security research and sci-fi"),
    ("Coffee",   "Always hunting for the best local roasters"),
    ("Training", "Weightlifting and outdoor runs"),
    ("Security", "CTFs, cloud security labs, homelab builds"),
]
cw = Inches(2.98); ch = Inches(1.85); gap = Inches(0.13)
for i, (title, body) in enumerate(cards):
    l = Inches(0.75) + i * (cw + gap)
    grid_cell_dark(sl, i+1, title, body, l, Inches(3.1), cw, ch)

# ═══════════════════════════════════════════════════════════════════
# S3 — THE SCENARIO (white)
# ═══════════════════════════════════════════════════════════════════
sl = slide(WHITE)
label_txt(sl, "01  The Scenario", dark=False)
rule_line(sl, Inches(0.88), dark=False)
h1(sl, "Meet Helix Labs", Inches(0.98), size=Pt(42), dark=False)

# Scenario bar
rect(sl, Inches(0.75), Inches(1.95), Inches(11.83), Inches(0.85),
     fill=OFFWHITE, border=RGBColor(0xe0,0xe0,0xdd))
tb(sl,
   "A 200-person biotech company. Two offices in San Francisco. Verkada cameras at every entrance, lab, "
   "and parking structure. Three-person IT team managing everything remotely.",
   Inches(0.98), Inches(2.06), Inches(11.3), Inches(0.72),
   size=Pt(13), color=BLACK, wrap=True)

# Column headers
for lx, text in [(Inches(0.75), "FUNCTIONAL REQUIREMENTS"), (Inches(6.9), "NON-FUNCTIONAL REQUIREMENTS")]:
    tb(sl, text, lx, Inches(3.0), Inches(5.8), Inches(0.26),
       size=Pt(8), bold=True, color=RGBColor(0x99,0x99,0x99))
    tb(sl, "What the system must do" if "FUNC" in text else "How the system must perform",
       lx, Inches(3.28), Inches(5.8), Inches(0.24), size=Pt(10), color=RGBColor(0xbb,0xbb,0xbb))

fr = [
    "Isolate camera systems from corporate LAN — no shared broadcast domain",
    "Single authenticated, auditable remote access point for IT admins",
    "Encrypted footage in private cloud — inaccessible from the public internet",
    "Centralized management of all Verkada cameras across both office locations",
    "Deploy and tear down the full environment in under two minutes via IaC",
]
nfr = [
    "Cost efficiency — pay-as-you-go cloud spend; no dedicated hardware CapEx",
    "Ease of maintenance — version-controlled IaC with zero configuration drift",
    "Network isolation — enforced by architecture, not configuration alone",
    "Auditability & reporting — all access logged; SOC 2, ISO 27001, HIPAA compliant",
    "Multi-client scalability — architecture templates across future client deployments",
]
item_list(sl, fr,  Inches(0.75), Inches(3.6), Inches(5.8), Inches(3.6), dark=False,
          num_col=RGBColor(0xcc,0xcc,0xcc))
item_list(sl, nfr, Inches(6.9),  Inches(3.6), Inches(6.1), Inches(3.6), dark=False,
          num_col=RGBColor(0xcc,0xcc,0xcc))

# ═══════════════════════════════════════════════════════════════════
# S4 — THE THREAT (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "01  The Scenario")
rule_line(sl, Inches(0.88))
h1(sl, "When every device shares the same network,\none breach can compromise everything.",
   Inches(0.98), size=Pt(38),
   bold_parts=["one breach can compromise everything."])

examples = [
    ("2021 — Verkada",
     "Single admin credential compromised. 150,000 cameras exposed — hospitals, jails, factories. One flat network."),
    ("2022 — Uber",
     "Contractor credentials stolen. One compromised account became a bridge to every internal system."),
    ("2023 — MGM Resorts",
     "Ten-minute phone call to IT helpdesk. One employee social-engineered. $100M in damages. Same pattern."),
]
ew = Inches(3.88); gap = Inches(0.175)
for i, (title, body) in enumerate(examples):
    l = Inches(0.75) + i * (ew + gap)
    rect(sl, l, Inches(3.45), ew, Inches(2.7),
         fill=RGBColor(0x12,0x04,0x04), border=RGBColor(0x45,0x12,0x12))
    # Red accent top bar
    rect(sl, l, Inches(3.45), ew, Pt(3), fill=RED)
    tb(sl, title, l+Inches(0.22), Inches(3.65), ew-Inches(0.3), Inches(0.34),
       size=Pt(12), bold=True, color=RED)
    tb(sl, body,  l+Inches(0.22), Inches(4.08), ew-Inches(0.3), Inches(1.9),
       size=Pt(13), color=LGRAY, wrap=True)

tb(sl,
   "The pattern is always the same: one weak link in a flat network gives an attacker lateral access to everything.",
   Inches(0.75), Inches(6.42), Inches(11.5), Inches(0.5),
   size=Pt(14), color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# S5 — SOLUTION CONSIDERATIONS (white)
# ═══════════════════════════════════════════════════════════════════
sl = slide(WHITE)
label_txt(sl, "02  The Solution", dark=False)
rule_line(sl, Inches(0.88), dark=False)
h1(sl, "Before building anything, we evaluated the tradeoffs.",
   Inches(0.98), size=Pt(38), dark=False)

cells = [
    ("Cost",                   "AWS pay-as-you-go ~$45/month. On-prem NVR starts at $50K+ CapEx. 60–70% lower TCO over 3 years."),
    ("Ease of Maintenance",    "Every resource is Terraform code. Any engineer can rebuild the full environment in one command. No undocumented state."),
    ("Network Isolation",      "Private subnet has no Internet Gateway route — no path exists by design, not by firewall rule. Isolation is structural."),
    ("Total Cost of Ownership","Managed AWS absorbs patching, hardware failure, and scaling. Helix Labs pays for what they use and nothing more."),
    ("Reporting & Auditability","CloudTrail, VPC Flow Logs, S3 access logging built in. SOC 2, ISO 27001, HIPAA satisfied out of the box."),
    ("Multi-Client Scalability","Terraform module is parameterized. Deploy the same zero-trust architecture for any client in under two minutes."),
]
cw = Inches(3.98); ch = Inches(2.0); gap = Inches(0.1)
for i, (title, body) in enumerate(cells):
    col = i % 3; row = i // 3
    l = Inches(0.75) + col * (cw + gap)
    t = Inches(2.3)  + row * (ch + gap)
    grid_cell_light(sl, i+1, title, body, l, t, cw, ch)

# ═══════════════════════════════════════════════════════════════════
# S6 — THE SOLUTION: SIX LAYERS (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "02  The Solution")
rule_line(sl, Inches(0.88))
h1(sl, "Six layers of protection. One secure environment.", Inches(0.98), size=Pt(38))

layers = [
    ("VPC",               "Helix Labs' own private section of the cloud. No other tenant can see in or route to it."),
    ("Subnet Segmentation","Public zone for admin access. Private zone for cameras. No shared broadcast domain."),
    ("Bastion Host",       "The only door in. SSH locked to one authorized IP. Every session is logged."),
    ("Security Groups",    "Deny-all by default. Camera server can upload footage and nothing else. Cannot receive connections."),
    ("S3 + IAM",           "Encrypted storage. Three actions, one bucket, no static credentials stored anywhere."),
    ("Terraform IaC",      "Everything defined as code. Peer-reviewable, reproducible, deployed in under two minutes."),
]
cw = Inches(3.98); ch = Inches(1.92); gap = Inches(0.1)
for i, (title, body) in enumerate(layers):
    col = i % 3; row = i // 3
    l = Inches(0.75) + col * (cw + gap)
    t = Inches(2.25) + row * (ch + gap)
    grid_cell_dark(sl, i+1, title, body, l, t, cw, ch)

# ═══════════════════════════════════════════════════════════════════
# S7 — ARCHITECTURE DIAGRAM (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "02  The Solution")
rule_line(sl, Inches(0.88))
h1(sl, "Physical architecture", Inches(0.98), size=Pt(38))

# VPC outer box
rect(sl, Inches(1.5), Inches(2.3), Inches(10.1), Inches(4.55),
     fill=RGBColor(0x0c,0x0c,0x10), border=RGBColor(0x28,0x28,0x32), bw=Pt(1))
tb(sl, "AWS VPC  ·  10.0.0.0/16", Inches(1.62), Inches(2.38), Inches(5), Inches(0.28),
   size=Pt(8), color=RGBColor(0x38,0x38,0x45))

# Public subnet
rect(sl, Inches(1.8), Inches(2.65), Inches(2.35), Inches(3.85),
     fill=RGBColor(0x08,0x12,0x0a), border=RGBColor(0x14,0x30,0x1a), bw=Pt(0.75))
tb(sl, "PUBLIC  ·  10.0.1.0/24", Inches(1.92), Inches(2.72), Inches(2.1), Inches(0.25),
   size=Pt(7.5), color=RGBColor(0x20,0x55,0x28))

# Private subnet
rect(sl, Inches(4.3), Inches(2.65), Inches(7.15), Inches(3.85),
     fill=RGBColor(0x08,0x0c,0x16), border=RGBColor(0x14,0x1e,0x38), bw=Pt(0.75))
tb(sl, "PRIVATE  ·  10.0.2.0/24", Inches(4.42), Inches(2.72), Inches(3), Inches(0.25),
   size=Pt(7.5), color=RGBColor(0x20,0x28,0x55))

def node(sl, l, t, w, h, title, subtitle, accent):
    rect(sl, l, t, w, h, fill=RGBColor(0x14,0x14,0x18), border=accent, bw=Pt(1.2), rounding=True)
    tb(sl, title,    l+Inches(0.08), t+Inches(0.1),  w-Inches(0.12), Inches(0.32),
       size=Pt(11), bold=True, color=accent, align=PP_ALIGN.CENTER)
    tb(sl, subtitle, l+Inches(0.08), t+Inches(0.44), w-Inches(0.12), Inches(0.25),
       size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER)

def arr_lbl(sl, text, l, t, w, color):
    tb(sl, text, l, t, w, Inches(0.26), size=Pt(9), color=color, align=PP_ALIGN.CENTER)

NY = Inches(3.85)
node(sl, Inches(0.08), NY, Inches(1.28), Inches(0.82), "INTERNET",   "0.0.0.0/0",          MUTED)
node(sl, Inches(1.98), NY, Inches(1.98), Inches(0.82), "BASTION",    "port 22 · CIDR-locked", GREEN)
node(sl, Inches(4.48), NY, Inches(1.92), Inches(0.82), "CAMERAS",    "Verkada · RTSP",       CYAN)
node(sl, Inches(6.6),  NY, Inches(2.08), Inches(0.82), "CAM SERVER", "NVR · processes feeds", BLUE)
node(sl, Inches(8.9),  NY, Inches(1.92), Inches(0.82), "S3 BUCKET",  "AES-256 · Block Public",AMBER)
node(sl, Inches(11.05),NY, Inches(1.92), Inches(0.82), "VIEWER",     "Admin / Client",        LGRAY)

# Arrow labels above nodes
arr_lbl(sl, "SSH :22 · via IGW",      Inches(0.0),  Inches(3.55), Inches(2.5),  GREEN)
arr_lbl(sl, "SSH proxy :22 (admin)",  Inches(3.8),  Inches(3.1),  Inches(2.8),  GREEN)
arr_lbl(sl, "RTSP stream",            Inches(4.5),  Inches(4.82), Inches(2.0),  CYAN)
arr_lbl(sl, "VPC EP :443 (write)",    Inches(6.55), Inches(4.82), Inches(2.3),  AMBER)
arr_lbl(sl, "HTTPS :443 · Signed URL",Inches(8.85), Inches(4.82), Inches(3.0),  LGRAY)

# ═══════════════════════════════════════════════════════════════════
# S8 — NETWORK SEGMENTATION (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "02  The Solution")
rule_line(sl, Inches(0.88))
h1(sl, "Network segmentation by architecture", Inches(0.98), size=Pt(38),
   bold_parts=["by architecture"])

items_net = [
    "A /16 VPC gives Helix Labs a dedicated address space — fully isolated from every other AWS tenant by default",
    "Public subnet (/24) routes through an Internet Gateway for admin access. Private subnet has no IGW route entry.",
    "The private subnet is unreachable by architecture — no route table entry means no path exists, not just a firewall rule blocking one",
    "All ingress to private hosts is SSH-proxied through the bastion host in the public subnet via agent forwarding",
]
item_list(sl, items_net, Inches(0.75), Inches(2.55), Inches(7.2), Inches(4.4))

callout(sl, "Why this matters",
    "Firewall rules can be misconfigured. Route tables cannot route to an address that has no entry. "
    "The private subnet is unreachable by design, not policy.",
    Inches(8.3), Inches(2.55), Inches(4.75), Inches(2.5))

# ═══════════════════════════════════════════════════════════════════
# S9 — BASTION HOST (white break)
# ═══════════════════════════════════════════════════════════════════
sl = slide(WHITE)
label_txt(sl, "02  The Solution", dark=False)
rule_line(sl, Inches(0.88), dark=False)
h1(sl, "Bastion host. One CIDR-locked ingress.\nNo lateral movement possible.",
   Inches(0.98), size=Pt(42), dark=False, bold_parts=["One CIDR-locked ingress."])

tb(sl,
   "Every admin session goes through one door. That door is locked to one IP address. "
   "Every connection attempt, every command, every session is logged.\n\n"
   "If an attacker reaches the camera server from inside — they still can't get out. "
   "The bastion doesn't let them back through. They're in the vault with no exit.",
   Inches(0.75), Inches(3.62), Inches(9.5), Inches(2.8),
   size=Pt(19), color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# S10 — SECURITY GROUPS (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "02  The Solution")
rule_line(sl, Inches(0.88))
h1(sl, "Zero inbound. Port 443 egress only.\nDeny-all default posture.", Inches(0.98), size=Pt(38))

items_sg = [
    "No inbound rules from the public internet — the camera server cannot be reached directly from anywhere.",
    "Ingress restricted to SSH from the bastion security group ID only — not an IP range, a specific group.",
    "Egress limited to port 443 HTTPS only — the server can upload footage to S3 and nothing else.",
    "AWS security groups are implicit deny-all by default. Every permission is explicitly defined. Nothing slips through.",
]
item_list(sl, items_sg, Inches(0.75), Inches(2.55), Inches(7.2), Inches(4.4))

callout(sl, "Think of it this way",
    "The camera server can make a phone call out, but it cannot receive one. "
    "Even if someone knew exactly where it was, they could not knock on the door.",
    Inches(8.3), Inches(2.55), Inches(4.75), Inches(2.5))

# ═══════════════════════════════════════════════════════════════════
# S11 — IAM (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "02  The Solution")
rule_line(sl, Inches(0.88))
h1(sl, "Least-privilege IAM.\nNo static credentials. Ever.", Inches(0.98), size=Pt(38))

items_iam = [
    "S3 bucket with Block Public Access enforced at account level — no accidental exposure, regardless of any other setting.",
    "EC2 instance role — temporary credentials issued automatically by AWS. No username or password stored anywhere.",
    "IAM policy grants exactly three actions: put a file, get a file, list the bucket. Nothing else is permitted.",
    "Policy scoped to one bucket ARN. Wildcard permissions explicitly denied. Least privilege enforced in code.",
]
item_list(sl, items_iam, Inches(0.75), Inches(2.55), Inches(7.2), Inches(4.4))

callout(sl, "Think of it this way",
    "The server has a badge that only opens one door and only lets it do three things once inside. "
    "Even if that badge were copied, it would be useless everywhere else.",
    Inches(8.3), Inches(2.55), Inches(4.75), Inches(2.5))

# ═══════════════════════════════════════════════════════════════════
# S12 — THE INCIDENT (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "03  The Moment It Matters")
rule_line(sl, Inches(0.88))
h1(sl, "Someone walks in who should not be there.", Inches(0.98), size=Pt(40))

# scenario box
rect(sl, Inches(0.75), Inches(2.6), Inches(11.83), Inches(1.05),
     fill=MGRAY, border=BORDER)
tb(sl,
   "It is 11:47 PM. An unrecognized person enters through the main entrance at Helix Labs. "
   "No badge. No record in the employee database. Confidence score: 94%.",
   Inches(0.98), Inches(2.72), Inches(11.3), Inches(0.85),
   size=Pt(15), color=LGRAY, wrap=True)

# alert card
rect(sl, Inches(0.75), Inches(3.88), Inches(11.83), Inches(1.72),
     fill=RGBColor(0x12,0x04,0x04), border=RGBColor(0x50,0x12,0x12))
rect(sl, Inches(0.75), Inches(3.88), Inches(11.83), Pt(3), fill=RED)
tb(sl, "⚠   UNAUTHORIZED ACCESS DETECTED",
   Inches(0.98), Inches(4.02), Inches(9), Inches(0.36),
   size=Pt(12), bold=True, color=RED)
tb(sl, "CAM-01  ·  Main Entrance  ·  HQ North  ·  No match in database  ·  Alert sent to security team",
   Inches(0.98), Inches(4.48), Inches(11.2), Inches(0.95),
   size=Pt(15), color=WHITE, wrap=True)

tb(sl,
   "The system doesn't wait for a human to notice. Alert fires automatically — security team notified, "
   "event logged with camera, timestamp, zone, and confidence score.",
   Inches(0.75), Inches(5.85), Inches(11.5), Inches(0.75),
   size=Pt(14), color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# S13 — WHAT HELIX LABS SEES (black)
# ═══════════════════════════════════════════════════════════════════
sl = slide(BLACK)
label_txt(sl, "03  The Moment It Matters")
rule_line(sl, Inches(0.88))
h1(sl, "What Helix Labs sees in real time.", Inches(0.98), size=Pt(40))

panels = [
    ("Right Now",   "Live Occupancy",       "How many people are in frame, identified vs. unknown, updated every second."),
    ("Who Is Here", "Employee Recognition", "Name, department, confidence score, and time on site."),
    ("Time Data",   "Dwell Time",           "How long each person has been in the space. Flags anyone staying unusually long."),
    ("Security",    "Unauthorized Alerts",  "One click to notify security, add to watchlist, or pull the full incident report."),
]
pw = Inches(2.98); ph = Inches(3.0); gap = Inches(0.12)
for i, (label_t, mid_t, body) in enumerate(panels):
    l = Inches(0.75) + i * (pw + gap)
    rect(sl, l, Inches(2.55), pw, ph, fill=MGRAY, border=BORDER)
    # Top accent
    rect(sl, l, Inches(2.55), pw, Pt(3), fill=BLUE)
    tb(sl, label_t, l+Inches(0.22), Inches(2.72), pw-Inches(0.3), Inches(0.28),
       size=Pt(9), color=MUTED)
    tb(sl, mid_t,   l+Inches(0.22), Inches(3.08), pw-Inches(0.3), Inches(0.38),
       size=Pt(15), bold=True, color=WHITE, wrap=True)
    tb(sl, body,    l+Inches(0.22), Inches(3.58), pw-Inches(0.3), ph-Inches(1.18),
       size=Pt(13), color=LGRAY, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# S14 — DEPLOYMENT (white break)
# ═══════════════════════════════════════════════════════════════════
sl = slide(WHITE)
label_txt(sl, "04  Deployment", dark=False)
rule_line(sl, Inches(0.88), dark=False)
h1(sl, "The entire system. Live in under two minutes.\nTorn down just as fast.",
   Inches(0.98), size=Pt(42), dark=False,
   bold_parts=["Live in under two minutes."])

tb(sl,
   "One command. The entire environment — private network, both zones, security desk, "
   "firewall rules, storage, permissions — built and running.\n\n"
   "No undocumented configuration. No single engineer who's the only one who knows how it works. "
   "It is a file in version control. Any engineer can read it, review it, and reproduce it exactly.",
   Inches(0.75), Inches(3.65), Inches(9.5), Inches(3.1),
   size=Pt(19), color=MUTED, wrap=True)

# ═══════════════════════════════════════════════════════════════════
# S15 — CONCLUSION (white)
# ═══════════════════════════════════════════════════════════════════
sl = slide(WHITE)
label_txt(sl, "05  Conclusion", dark=False)
rule_line(sl, Inches(0.88), dark=False)
h1(sl, "Requirements met. Principles upheld.", Inches(0.98), size=Pt(38), dark=False)

# ── Left: considerations + tenets ─────────────────────────────────
tb(sl, "SOLUTION CONSIDERATIONS", Inches(0.75), Inches(2.5), Inches(5.6), Inches(0.26),
   size=Pt(8), bold=True, color=RGBColor(0x99,0x99,0x99))

rows = [
    ("COST",   "$45/mo cloud vs $50K+ on-prem — 60% lower 3-yr TCO, zero CapEx"),
    ("MAINT",  "Terraform IaC — zero config drift, peer-reviewable, rebuild in one command"),
    ("ISOLN",  "Private subnet has no IGW route — isolated by architecture, not firewall rules"),
    ("SCALE",  "One Terraform module, parameterized — same deployment for every future client"),
]
for i, (tag, body) in enumerate(rows):
    t = Inches(2.82) + i * Inches(0.74)
    rect(sl, Inches(0.75), t, Inches(5.6), Inches(0.68),
         fill=OFFWHITE, border=RGBColor(0xe0,0xe0,0xdd))
    tb(sl, tag,  Inches(0.95), t+Inches(0.12), Inches(0.88), Inches(0.42),
       size=Pt(8), bold=True, color=BLACK)
    tb(sl, body, Inches(1.92), t+Inches(0.12), Inches(4.3),  Inches(0.58),
       size=Pt(11), color=MUTED, wrap=True)

tb(sl, "TENETS", Inches(0.75), Inches(5.84), Inches(5.6), Inches(0.26),
   size=Pt(8), bold=True, color=RGBColor(0x99,0x99,0x99))
tenets = [
    "01  Zero trust — no implicit access, every path explicitly permitted and logged",
    "02  Least privilege — three S3 actions, one bucket ARN, no static credentials",
    "03  Defense in depth — six independent layers, each limits blast radius",
]
for i, t_text in enumerate(tenets):
    tb(sl, t_text, Inches(0.75), Inches(6.12)+i*Inches(0.28), Inches(5.6), Inches(0.28),
       size=Pt(12), color=MUTED)

# ── Right: requirements + stats ────────────────────────────────────
tb(sl, "ALL REQUIREMENTS MET", Inches(6.75), Inches(2.5), Inches(6.3), Inches(0.26),
   size=Pt(8), bold=True, color=RGBColor(0x99,0x99,0x99))

reqs = [
    "FR — Camera systems isolated from corporate LAN by architecture",
    "FR — Single auditable ingress; every admin session is logged",
    "FR — Encrypted footage in private S3; unreachable without IAM auth",
    "NFR — Cost-justified: 60% lower TCO vs. on-prem over 3 years",
    "NFR — Multi-client: same Terraform module deploys for any client in <2 min",
    "NFR — Compliance-ready: CloudTrail + VPC Flow Logs cover SOC 2 & ISO 27001",
]
for i, req in enumerate(reqs):
    t = Inches(2.82) + i * Inches(0.58)
    # checkmark
    tb(sl, "✓", Inches(6.75), t, Inches(0.38), Inches(0.38),
       size=Pt(14), bold=True, color=GREEN)
    tb(sl, req, Inches(7.18), t, Inches(5.8), Inches(0.38),
       size=Pt(12), color=BLACK, wrap=True)

# Stats bar
stats = [("19","Cloud Resources"),("0","Exposed Entry Points"),("<2m","Deploy Time"),("6/6","Requirements Met")]
sw = Inches(1.46); gap_s = Inches(0.085)
for i, (n, lbl) in enumerate(stats):
    l = Inches(6.75) + i * (sw + gap_s)
    rect(sl, l, Inches(6.32), sw, Inches(0.92),
         fill=OFFWHITE, border=RGBColor(0xe0,0xe0,0xdd))
    tb(sl, n,   l+Inches(0.08), Inches(6.36), sw-Inches(0.1), Inches(0.44),
       size=Pt(24), color=BLACK, align=PP_ALIGN.CENTER)
    tb(sl, lbl, l+Inches(0.06), Inches(6.8),  sw-Inches(0.08),Inches(0.4),
       size=Pt(8), color=MUTED, align=PP_ALIGN.CENTER, wrap=True)

# ── Save ─────────────────────────────────────────────────────────
out = "/Users/samyu/verkada-infra/verkada-presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
